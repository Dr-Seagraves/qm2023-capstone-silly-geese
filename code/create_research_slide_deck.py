from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Color system: navy dominant, light blue support, bright blue accents.
NAVY = RGBColor(11, 31, 58)
NAVY_ALT = RGBColor(17, 45, 84)
LIGHT_BLUE = RGBColor(220, 235, 255)
PALE_BLUE = RGBColor(240, 246, 255)
ACCENT = RGBColor(92, 169, 255)
TEXT_DARK = RGBColor(22, 33, 52)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_header_band(slide, title, dark=False):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY_ALT if dark else LIGHT_BLUE
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(11.8), Inches(0.5))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.name = "Calibri"
    p.font.color.rgb = WHITE if dark else NAVY


def add_footer(slide, txt, dark=False):
    f = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.3))
    tf = f.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(11)
    p.font.name = "Calibri"
    p.font.color.rgb = LIGHT_BLUE if dark else NAVY_ALT


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    # Background motif blocks.
    block1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.8), Inches(4.3), Inches(2.1))
    block1.fill.solid()
    block1.fill.fore_color.rgb = NAVY_ALT
    block1.line.fill.background()

    block2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(3.15), Inches(4.9), Inches(2.5))
    block2.fill.solid()
    block2.fill.fore_color.rgb = ACCENT
    block2.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(7.0), Inches(2.8))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Lobbying and Firm Profitability"
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.name = "Cambria"
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "A Panel-Data Research Synthesis (2010-2020)"
    p2.font.size = Pt(22)
    p2.font.name = "Calibri"
    p2.font.color.rgb = LIGHT_BLUE

    meta = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8.6), Inches(1.0))
    mt = meta.text_frame
    mt.clear()
    m1 = mt.paragraphs[0]
    m1.text = "QM 2023 Capstone Team"
    m1.font.size = Pt(20)
    m1.font.bold = True
    m1.font.name = "Calibri"
    m1.font.color.rgb = WHITE

    m2 = mt.add_paragraph()
    m2.text = "Alycia Reji • Gracie Vivion • Shelby Howard • Daniz Mammadova"
    m2.font.size = Pt(15)
    m2.font.name = "Calibri"
    m2.font.color.rgb = LIGHT_BLUE

    add_footer(slide, "Research deck generated from project outputs in this repository", dark=True)


def add_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PALE_BLUE)
    add_header_band(slide, "Agenda")

    agenda_items = [
        "1. Research question and hypotheses",
        "2. Data pipeline, coverage, and quality",
        "3. EDA insights: correlation, lag structure, and heterogeneity",
        "4. Econometric design: two-way fixed effects + robustness",
        "5. Diagnostics, forecast benchmark, and interpretation",
        "6. Conclusions, limitations, and next steps",
    ]

    y = 1.35
    for item in agenda_items:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(11.9), Inches(0.78))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_BLUE

        t = slide.shapes.add_textbox(Inches(1.05), Inches(y + 0.18), Inches(11.2), Inches(0.38))
        p = t.text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.name = "Calibri"
        p.font.color.rgb = TEXT_DARK
        y += 0.9

    add_footer(slide, "Theme: midnight executive (navy + ice blue)")


def add_question_hypotheses_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_band(slide, "Research Question and Testable Hypotheses")

    q_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.15), Inches(12.0), Inches(1.05))
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = PALE_BLUE
    q_box.line.color.rgb = LIGHT_BLUE

    q_text = slide.shapes.add_textbox(Inches(1.05), Inches(1.38), Inches(11.5), Inches(0.6))
    q = q_text.text_frame.paragraphs[0]
    q.text = "What is the relationship between firm lobbying expenditures and subsequent profitability?"
    q.font.bold = True
    q.font.size = Pt(22)
    q.font.name = "Cambria"
    q.font.color.rgb = NAVY

    left = slide.shapes.add_textbox(Inches(0.9), Inches(2.45), Inches(5.95), Inches(4.3))
    ltf = left.text_frame
    ltf.clear()
    items_l = [
        "H1: Higher lobbying spend is associated with higher subsequent performance.",
        "H2: Effects appear with 1-2 year lags rather than contemporaneously.",
        "H3a: Effect magnitude differs by firm size.",
    ]
    for i, txt in enumerate(items_l):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = txt
        p.level = 0
        p.font.size = Pt(17)
        p.font.name = "Calibri"
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(12)

    right = slide.shapes.add_textbox(Inches(6.65), Inches(2.45), Inches(5.9), Inches(4.3))
    rtf = right.text_frame
    rtf.clear()
    items_r = [
        "H3b: Effects are stronger in policy-exposed industries.",
        "H4: Positive relationship should survive log-scaling/outlier robustness.",
        "H5: Controls should attenuate, not fully erase, the lobbying coefficient.",
    ]
    for i, txt in enumerate(items_r):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = txt
        p.level = 0
        p.font.size = Pt(17)
        p.font.name = "Calibri"
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(12)

    add_footer(slide, "Hypothesis set compiled from EDA and M3 interpretation docs")


def add_data_coverage_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PALE_BLUE)
    add_header_band(slide, "Data Pipeline and Coverage Snapshot")

    cards = [
        ("Financials", "4,932 rows\n1,375 firms\nSEC 10-K (2010-2020)"),
        ("Lobbying", "11,619 rows\n1,534 firms\nSenate filings"),
        ("Merged panel", "5,099 rows\n1,375 firms\nfirm-year panel"),
        ("Balanced panel", "836 rows\n66 firms\nstrict 2010-2020"),
    ]

    x = 0.8
    for title, body in cards:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.4), Inches(3.0), Inches(2.05))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = LIGHT_BLUE

        t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.62), Inches(2.6), Inches(0.5))
        p1 = t.text_frame.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(19)
        p1.font.name = "Cambria"
        p1.font.color.rgb = NAVY

        b = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.1), Inches(2.6), Inches(1.15))
        p2 = b.text_frame.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(14)
        p2.font.name = "Calibri"
        p2.font.color.rgb = TEXT_DARK
        x += 3.15

    bullets = slide.shapes.add_textbox(Inches(0.9), Inches(3.9), Inches(12.0), Inches(2.8))
    btf = bullets.text_frame
    btf.clear()
    notes = [
        "Crosswalk coverage is the key bottleneck: 1,428 rows map to gvkey; 3,671 do not.",
        "Only 245 merged rows contain non-missing lobbying spend (4.8% of merged panel).",
        "Balanced panel improves temporal consistency but removes 83.6% of rows.",
        "Interpret estimates with explicit attention to sparse treatment support.",
    ]
    for i, txt in enumerate(notes):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = txt
        p.level = 0
        p.font.size = Pt(16)
        p.font.name = "Calibri"
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    add_footer(slide, "Source: quality_report.md")


def add_data_quality_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_band(slide, "Data Quality Risks and Mitigations")

    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(6.1), Inches(5.7))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = PALE_BLUE
    left_card.line.color.rgb = LIGHT_BLUE

    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(1.35), Inches(5.55), Inches(5.7))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = LIGHT_BLUE

    l = slide.shapes.add_textbox(Inches(1.1), Inches(1.65), Inches(5.6), Inches(5.0)).text_frame
    l.clear()
    left_lines = [
        "Main Risks",
        "• Missing lobbying spend: 95.2% of merged rows",
        "• ROA distribution has extreme tails (mean unstable)",
        "• Crosswalk incompleteness limits merged support",
        "• Potential omitted-variable and timing bias",
        "",
        "Practical Implication",
        "Need robust inference, outlier checks, lag tests, and cautious causal language.",
    ]
    for i, txt in enumerate(left_lines):
        p = l.paragraphs[0] if i == 0 else l.add_paragraph()
        p.text = txt
        p.font.name = "Calibri"
        p.font.size = Pt(20 if i in [0, 6] else 16)
        p.font.bold = i in [0, 6]
        p.font.color.rgb = NAVY if i in [0, 6] else TEXT_DARK

    r = slide.shapes.add_textbox(Inches(7.2), Inches(1.65), Inches(5.05), Inches(5.0)).text_frame
    r.clear()
    right_lines = [
        "What We Did",
        "• Two-way fixed effects (firm + year)",
        "• Clustered SE at firm level",
        "• Lag structure comparison (t-1, t-2, t-3)",
        "• Placebo lead test for timing/reverse causality",
        "• Excluding 2020 shock-year robustness",
        "• Size-based heterogeneity splits",
        "• Bootstrap interval for lobbying effect",
    ]
    for i, txt in enumerate(right_lines):
        p = r.paragraphs[0] if i == 0 else r.add_paragraph()
        p.text = txt
        p.font.name = "Calibri"
        p.font.size = Pt(20 if i == 0 else 16)
        p.font.bold = i == 0
        p.font.color.rgb = NAVY if i == 0 else TEXT_DARK

    add_footer(slide, "Risk framing: association-first, causality guarded")


def add_eda_slide(prs, fig_root):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PALE_BLUE)
    add_header_band(slide, "EDA Signal: Positive Co-Movement, Moderate Strength")

    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(6.0), Inches(5.75))
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = LIGHT_BLUE

    heatmap = fig_root / "M2" / "plot1_correlation_heatmap.png"
    if heatmap.exists():
        slide.shapes.add_picture(str(heatmap), Inches(1.05), Inches(1.65), Inches(5.5), Inches(4.55))
    else:
        tx = slide.shapes.add_textbox(Inches(1.05), Inches(2.9), Inches(5.4), Inches(1.0))
        p = tx.text_frame.paragraphs[0]
        p.text = "Correlation heatmap unavailable"
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY

    right = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.7), Inches(5.3)).text_frame
    right.clear()
    lines = [
        "Evidence Summary",
        "• Lobbying and firm performance move together positively in raw data.",
        "• Relationship is moderate, not overwhelming.",
        "• Control additions attenuate coefficients but often preserve directional signal.",
        "• Outliers amplify magnitude in levels; sign tends to persist under robust checks.",
        "",
        "Interpretation",
        "Data is consistent with a selection-plus-impact narrative: stronger firms lobby more, and lobbying may add incremental association.",
    ]
    for i, txt in enumerate(lines):
        p = right.paragraphs[0] if i == 0 else right.add_paragraph()
        p.text = txt
        p.font.name = "Calibri"
        p.font.size = Pt(21 if i in [0, 6] else 16)
        p.font.bold = i in [0, 6]
        p.font.color.rgb = NAVY if i in [0, 6] else TEXT_DARK

    add_footer(slide, "Figure: M2 plot1_correlation_heatmap.png")


def add_lag_slide(prs, fig_root):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_band(slide, "Lag Dynamics: Effects Emerge Over Time")

    fig = fig_root / "M2" / "plot4_lagged_effect_analysis.png"
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(8.2), Inches(5.75))
    frame.fill.solid()
    frame.fill.fore_color.rgb = PALE_BLUE
    frame.line.color.rgb = LIGHT_BLUE

    if fig.exists():
        slide.shapes.add_picture(str(fig), Inches(1.05), Inches(1.7), Inches(7.7), Inches(4.9))

    txt = slide.shapes.add_textbox(Inches(9.2), Inches(1.5), Inches(3.95), Inches(5.3)).text_frame
    txt.clear()
    points = [
        "Model Evidence",
        "Lag 1: -174.6 (p=0.135, clustered)",
        "Lag 2: -95.3 (p=0.272)",
        "Lag 3: 38.5 (p=0.648)",
        "Lead 1 placebo: -302.7 (p=0.069)",
        "",
        "Takeaway",
        "Timing patterns are not cleanly causal; lead effect suggests potential anticipation or reverse-causality dynamics.",
    ]
    for i, txt_line in enumerate(points):
        p = txt.paragraphs[0] if i == 0 else txt.add_paragraph()
        p.text = txt_line
        p.font.name = "Calibri"
        p.font.size = Pt(20 if i in [0, 6] else 14)
        p.font.bold = i in [0, 6]
        p.font.color.rgb = NAVY if i in [0, 6] else TEXT_DARK

    add_footer(slide, "Lag and lead metrics from M3 robustness table")


def add_heterogeneity_slide(prs, fig_root):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PALE_BLUE)
    add_header_band(slide, "Heterogeneity and Leverage")

    fig1 = fig_root / "M2" / "plot8_outcome_by_driver_quintile.png"
    fig2 = fig_root / "M2" / "plot6_control_scatter_regression.png"

    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(6.1), Inches(3.0))
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = LIGHT_BLUE
    if fig1.exists():
        slide.shapes.add_picture(str(fig1), Inches(1.0), Inches(1.55), Inches(5.7), Inches(2.6))

    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.35), Inches(5.5), Inches(3.0))
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.color.rgb = LIGHT_BLUE
    if fig2.exists():
        slide.shapes.add_picture(str(fig2), Inches(7.2), Inches(1.55), Inches(5.1), Inches(2.6))

    txt = slide.shapes.add_textbox(Inches(0.9), Inches(4.55), Inches(11.8), Inches(2.2)).text_frame
    txt.clear()
    lines = [
        "• Group sensitivity: association appears stronger for larger and policy-exposed firms in EDA narrative.",
        "• FE heterogeneity split: small firms beta=-5753.9 (p=0.156, imprecise) vs large firms beta=-24.1 (p<0.001).",
        "• Outlier influence: extreme spenders/very large firms materially affect slope magnitude in levels.",
        "• Robust interpretation: concentration at the top changes effect size more than directional pattern in many checks.",
    ]
    for i, line in enumerate(lines):
        p = txt.paragraphs[0] if i == 0 else txt.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.name = "Calibri"
        p.font.color.rgb = TEXT_DARK

    add_footer(slide, "Figures: M2 plot8 and plot6")


def add_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_band(slide, "Econometric Design")

    eq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(12.0), Inches(1.35))
    eq.fill.solid()
    eq.fill.fore_color.rgb = PALE_BLUE
    eq.line.color.rgb = LIGHT_BLUE

    eq_txt = slide.shapes.add_textbox(Inches(1.05), Inches(1.8), Inches(11.4), Inches(0.6)).text_frame
    ep = eq_txt.paragraphs[0]
    ep.text = "ROA_it = beta*Lobbying_{i,t-1} + theta1*log(Assets_it) + theta2*log(Revenues_it) + alpha_i + gamma_t + epsilon_it"
    ep.font.name = "Consolas"
    ep.font.size = Pt(16)
    ep.font.color.rgb = NAVY

    columns = [
        ("Identification Features", [
            "Firm fixed effects remove time-invariant heterogeneity.",
            "Year fixed effects absorb macro shocks.",
            "Lagged lobbying reduces contemporaneous simultaneity.",
            "Clustered SEs handle within-firm dependence.",
        ]),
        ("Threats to Causal Interpretation", [
            "Lead placebo signal implies timing concern.",
            "Sparse non-missing lobbying reduces support.",
            "Unobserved time-varying strategy may remain.",
            "Results should be framed as association-first.",
        ]),
    ]

    xvals = [0.9, 6.7]
    for idx, (title, points) in enumerate(columns):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(xvals[idx]), Inches(3.0), Inches(5.7), Inches(3.8))
        card.fill.solid()
        card.fill.fore_color.rgb = PALE_BLUE if idx == 0 else WHITE
        card.line.color.rgb = LIGHT_BLUE

        tf = slide.shapes.add_textbox(Inches(xvals[idx] + 0.2), Inches(3.2), Inches(5.3), Inches(3.3)).text_frame
        tf.clear()
        for j, line in enumerate([title] + points):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = line if j == 0 else f"• {line}"
            p.font.size = Pt(19 if j == 0 else 14)
            p.font.name = "Calibri"
            p.font.bold = j == 0
            p.font.color.rgb = NAVY if j == 0 else TEXT_DARK

    add_footer(slide, "Model A: two-way FE panel with firm-clustered inference")


def add_results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PALE_BLUE)
    add_header_band(slide, "Main Results: FE Coefficients and Precision")

    table_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(8.5), Inches(5.7))
    table_card.fill.solid()
    table_card.fill.fore_color.rgb = WHITE
    table_card.line.color.rgb = LIGHT_BLUE

    rows = [
        ["Specification", "Lobbying Coef", "Std. Err", "p-value", "N"],
        ["Lag 1 FE (clustered)", "-174.6", "116.9", "0.135", "2125"],
        ["Lag 2 FE (clustered)", "-95.3", "86.7", "0.272", "1700"],
        ["Lag 3 FE (clustered)", "38.5", "84.3", "0.648", "1372"],
        ["Exclude 2020", "-109.2", "72.2", "0.131", "1911"],
    ]

    t = slide.shapes.add_table(len(rows), len(rows[0]), Inches(1.1), Inches(1.8), Inches(7.9), Inches(4.6)).table
    col_widths = [Inches(2.8), Inches(1.35), Inches(1.25), Inches(1.15), Inches(0.9)]
    for i, w in enumerate(col_widths):
        t.columns[i].width = w

    for r in range(len(rows)):
        for c in range(len(rows[0])):
            cell = t.cell(r, c)
            cell.text = rows[r][c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE if r == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Calibri"
            p.font.size = Pt(13)
            p.font.bold = r == 0
            p.font.color.rgb = NAVY if r == 0 else TEXT_DARK
            if c > 0:
                p.alignment = PP_ALIGN.CENTER

    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.55), Inches(1.35), Inches(3.0), Inches(5.7))
    insight.fill.solid()
    insight.fill.fore_color.rgb = NAVY
    insight.line.fill.background()

    tf = slide.shapes.add_textbox(Inches(9.8), Inches(1.75), Inches(2.5), Inches(4.9)).text_frame
    tf.clear()
    lines = [
        "Interpretation",
        "Point estimates are often negative, but statistical precision is limited in most specifications.",
        "The lead-placebo and instability across lags reinforce an association-based, non-causal reading.",
        "Coefficient movement with specification changes is economically meaningful.",
    ]
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.name = "Calibri"
        p.font.size = Pt(18 if i == 0 else 13)
        p.font.bold = i == 0
        p.font.color.rgb = WHITE if i == 0 else LIGHT_BLUE
        p.space_after = Pt(8)

    add_footer(slide, "Tabulated from M3 robustness outputs")


def add_diagnostics_slide(prs, fig_root):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_band(slide, "Diagnostics and Inference")

    fig_qq = fig_root / "M3_qq_plot.png"
    fig_rf = fig_root / "M3_residuals_vs_fitted.png"

    c1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(6.2), Inches(3.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = PALE_BLUE
    c1.line.color.rgb = LIGHT_BLUE
    if fig_qq.exists():
        slide.shapes.add_picture(str(fig_qq), Inches(1.0), Inches(1.55), Inches(5.8), Inches(2.6))

    c2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.35), Inches(5.5), Inches(3.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = PALE_BLUE
    c2.line.color.rgb = LIGHT_BLUE
    if fig_rf.exists():
        slide.shapes.add_picture(str(fig_rf), Inches(7.2), Inches(1.55), Inches(5.1), Inches(2.6))

    tf = slide.shapes.add_textbox(Inches(0.9), Inches(4.6), Inches(11.9), Inches(2.1)).text_frame
    tf.clear()
    notes = [
        "• Breusch-Pagan p-value = 0.0000: heteroskedasticity present, so robust/clustered inference is required.",
        "• Max VIF = 2.74: no serious multicollinearity among key regressors.",
        "• Residual diagnostics indicate non-ideal shape, consistent with sparse and heterogeneous panel structure.",
        "• Conclusion: standard FE estimates are informative, but robust uncertainty treatment is essential.",
    ]
    for i, n in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = n
        p.font.name = "Calibri"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_DARK

    add_footer(slide, "Figures: M3 diagnostic plots and diagnostics summary")


def add_arima_slide(prs, fig_root):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PALE_BLUE)
    add_header_band(slide, "Forecast Benchmark (ARIMA)")

    fig = fig_root / "M3_arima_forecast.png"
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(8.2), Inches(5.7))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_BLUE
    if fig.exists():
        slide.shapes.add_picture(str(fig), Inches(1.1), Inches(1.7), Inches(7.6), Inches(4.95))

    tf = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.35), Inches(3.35), Inches(5.7))
    tf.fill.solid()
    tf.fill.fore_color.rgb = NAVY
    tf.line.fill.background()

    txt = slide.shapes.add_textbox(Inches(9.45), Inches(1.65), Inches(2.9), Inches(5.2)).text_frame
    txt.clear()
    lines = [
        "ARIMA Detail",
        "Order: (0,1,0)",
        "ADF p-value: 0.229",
        "ARIMA RMSE: 523.53",
        "Naive RMSE: 523.53",
        "",
        "Result",
        "ARIMA does not outperform naive baseline in this annual sample.",
    ]
    for i, ln in enumerate(lines):
        p = txt.paragraphs[0] if i == 0 else txt.add_paragraph()
        p.text = ln
        p.font.name = "Calibri"
        p.font.size = Pt(18 if i in [0, 6] else 14)
        p.font.bold = i in [0, 6]
        p.font.color.rgb = WHITE if i in [0, 6] else LIGHT_BLUE

    add_footer(slide, "ARIMA included as predictive benchmark, not identification strategy")


def add_conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_header_band(slide, "Conclusions and Strategic Next Steps", dark=True)

    concl = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(12.1), Inches(4.6)).text_frame
    concl.clear()
    points = [
        "Key Conclusions",
        "• Evidence supports a meaningful lobbying-profitability association, but not a clean causal claim.",
        "• Coefficients are sensitive to timing and sample support, with lead effects raising endogeneity concerns.",
        "• Diagnostic and robustness layers improve credibility, yet sparse lobbying coverage remains the core constraint.",
        "",
        "Next Steps",
        "• Improve identifier mapping and expand policy-exposure controls.",
        "• Test alternative outcomes (ROE, margins, valuation) and dynamic panel estimators.",
        "• Explore quasi-experimental variation for stronger causal leverage.",
    ]
    for i, line in enumerate(points):
        p = concl.paragraphs[0] if i == 0 else concl.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.font.size = Pt(30 if i in [0, 5] else 18)
        p.font.bold = i in [0, 5]
        p.font.color.rgb = WHITE if i in [0, 5] else LIGHT_BLUE
        p.space_after = Pt(7)

    sig = slide.shapes.add_textbox(Inches(0.9), Inches(6.25), Inches(12.0), Inches(0.6))
    sp = sig.text_frame.paragraphs[0]
    sp.text = "Thank you"
    sp.font.name = "Cambria"
    sp.font.bold = True
    sp.font.size = Pt(26)
    sp.font.color.rgb = ACCENT
    sp.alignment = PP_ALIGN.RIGHT

    add_footer(slide, "QM 2023 Capstone | Lobbying and Firm Profitability", dark=True)


def build_deck(output_path: Path, fig_root: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_agenda_slide(prs)
    add_question_hypotheses_slide(prs)
    add_data_coverage_slide(prs)
    add_data_quality_slide(prs)
    add_eda_slide(prs, fig_root)
    add_lag_slide(prs, fig_root)
    add_heterogeneity_slide(prs, fig_root)
    add_model_slide(prs)
    add_results_slide(prs)
    add_diagnostics_slide(prs, fig_root)
    add_arima_slide(prs, fig_root)
    add_conclusion_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main():
    repo_root = Path(__file__).resolve().parents[1]
    fig_root = repo_root / "results" / "figures"
    output = repo_root / "results" / "reports" / "QM2023_Lobbying_Research_Deck.pptx"
    build_deck(output, fig_root)
    print(f"Created: {output}")


if __name__ == "__main__":
    main()